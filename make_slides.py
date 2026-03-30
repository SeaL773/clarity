"""Generate Clarity pitch deck as PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Colors
BG = RGBColor(0x0D, 0x15, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xB0, 0xB8, 0xC8)
MUTED = RGBColor(0x6B, 0x80, 0x90)
BLUE = RGBColor(0x5B, 0x7F, 0xBF)
ACCENT = RGBColor(0x8A, 0xB4, 0xF8)
ORANGE = RGBColor(0xE8, 0xA8, 0x70)
GREEN = RGBColor(0x66, 0xBB, 0x6A)
RED = RGBColor(0xE5, 0x73, 0x73)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=24, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_multi(slide, left, top, width, height, lines):
    """lines = [(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(8)

# ─── Slide 1: Title ───
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s)
add_text(s, 0, 1.5, 13.333, 1, "✨", 72, WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0, 2.5, 13.333, 1.2, "Clarity", 64, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 3.8, 10, 1, "AI-powered smart todo list for people who think\nfaster than they organize", 28, MUTED, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 5.2, 10, 0.5, "Flutter  •  FastAPI  •  Claude Haiku 4.5  •  Kiro IDE", 18, ACCENT, align=PP_ALIGN.CENTER)
add_text(s, 0, 6.5, 13.333, 0.5, "AWS Kiro × CS Careers Hackathon — Virginia Tech — March 2026", 16, MUTED, align=PP_ALIGN.CENTER)

# ─── Slide 2: Problem ───
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 1.2, 11, 1.5, "366M+", 96, BLUE, bold=True)
add_text(s, 1, 2.8, 11, 1, "people worldwide have ADHD", 48, WHITE, bold=True)
add_text(s, 1, 4.2, 9, 1.5, "The hardest part isn't doing the work — it's figuring out where to start.\n\nTraditional todo apps require you to already be organized. But what if you can't?", 24, LIGHT)
add_text(s, 1, 6.2, 10, 0.5, "Executive Dysfunction  •  Task Initiation  •  Planning Overhead", 18, RED)

# ─── Slide 3: Solution ───
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 1.5, 11, 1.5, "Just dump your thoughts.\nAI handles the rest.", 48, WHITE, bold=True)
add_text(s, 1, 3.8, 10, 1.2, '"I have a midterm Thursday, haven\'t started studying, need to do laundry, my room\'s a mess, and I told Sarah I\'d review her resume."', 22, MUTED)
add_text(s, 1, 5.3, 10, 0.5, "↓  One button  ↓", 28, ACCENT, align=PP_ALIGN.LEFT)
add_text(s, 1, 6.0, 10, 0.5, "Structured tasks  •  Prioritized  •  Sub-steps  •  Ready to go", 22, LIGHT)

# ─── Slide 4: AI Pipeline ───
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.8, 11, 1, "3-Step AI Pipeline", 48, WHITE, bold=True)
add_text(s, 1, 1.8, 10, 0.5, "Not a single API call — three dedicated prompts, each designed for its task.", 20, MUTED)
add_multi(s, 1, 2.8, 10, 1.2, [
    ("1. Parse → Extract Tasks", 28, ACCENT, True),
    ('Finds every task — even implicit ones. "I need to eat healthier" → "Plan meals" + "Go grocery shopping"', 20, LIGHT, False),
])
add_multi(s, 1, 4.2, 10, 1.2, [
    ("2. Plan → Prioritize & Order", 28, ORANGE, True),
    ("Eisenhower Matrix priority. Quick wins first for momentum. Workflow-aware ordering.", 20, LIGHT, False),
])
add_multi(s, 1, 5.6, 10, 1.2, [
    ("3. Summarize → Daily Recap", 28, GREEN, True),
    ("Encouraging summary. Leads with wins, never guilt. Suggests tomorrow's focus.", 20, LIGHT, False),
])

# ─── Slide 5: Features ───
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.8, 11, 1, "ADHD-Friendly by Design", 48, WHITE, bold=True)
# Row 1
add_multi(s, 1, 2.2, 5, 1.8, [
    ("🚫⏱️  No Time Pressure", 24, ACCENT, True),
    ("Zero timers, zero estimates. They cause anxiety — we removed them entirely.", 18, LIGHT, False),
])
add_multi(s, 7, 2.2, 5.5, 1.8, [
    ("🎉  Celebration Moments", 24, ACCENT, True),
    ("⚡ First done · 🔥 Halfway · 🎉 All complete.\nDopamine hits built into the UX.", 18, LIGHT, False),
])
# Row 2
add_multi(s, 1, 4.3, 5, 1.8, [
    ("📊  Smart Sorting", 24, ACCENT, True),
    ("Completed tasks sink. Priority colors guide you. Sub-tasks auto-complete parents.", 18, LIGHT, False),
])
add_multi(s, 7, 4.3, 5.5, 1.8, [
    ("🌙  Calm UI", 24, ACCENT, True),
    ("Warm dark mode. Frosted glass nav. Generous whitespace. Reduces cognitive load.", 18, LIGHT, False),
])

# ─── Slide 6: Calendar ───
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.8, 11, 1, "Calendar + Timeline", 48, WHITE, bold=True)
add_multi(s, 1, 2.2, 7, 4.5, [
    ("→  Monthly view with urgency progress bars per day", 22, LIGHT, False),
    ("→  Color depth = urgency (blue light / amber dark mode)", 22, LIGHT, False),
    ("→  Fill = completion percentage", 22, LIGHT, False),
    ("→  Google Calendar-style timeline with due times", 22, LIGHT, False),
    ("→  Tap to preview, double-tap to expand", 22, LIGHT, False),
    ("→  Quick 'Today' button to jump back", 22, LIGHT, False),
    ("→  Historical read-only view for past days", 22, LIGHT, False),
])
add_text(s, 9, 3, 3.5, 3, "📅", 120, WHITE, align=PP_ALIGN.CENTER)

# ─── Slide 7: Voice ───
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.8, 11, 1, "Voice Input", 48, WHITE, bold=True)
add_text(s, 1, 1.8, 10, 0.5, "Because sometimes typing feels like too much.", 24, MUTED)
add_multi(s, 1, 3, 5, 2.5, [
    ("On-Device STT", 28, ORANGE, True),
    ("Instant. Private. Works offline.", 20, LIGHT, False),
    ("Uses native speech_to_text.", 20, LIGHT, False),
])
add_multi(s, 7, 3, 5.5, 2.5, [
    ("OpenAI Whisper", 28, ACCENT, True),
    ("Record → Upload → Transcribe.", 20, LIGHT, False),
    ("Higher accuracy. Server-side.", 20, LIGHT, False),
])
add_text(s, 1, 5.8, 10, 0.5, "Switchable modes  •  Platform-aware  •  Bundled test audio for verification", 18, MUTED)

# ─── Slide 8: Tech Stack ───
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.8, 11, 1, "Tech Stack", 48, WHITE, bold=True)
add_multi(s, 1, 2.2, 3.5, 3, [
    ("Frontend", 28, RGBColor(0x02, 0x56, 0x9B), True),
    ("Flutter + Dart", 20, LIGHT, False),
    ("Material 3", 20, LIGHT, False),
    ("Provider state", 20, LIGHT, False),
    ("table_calendar", 20, LIGHT, False),
])
add_multi(s, 5, 2.2, 3.5, 3, [
    ("Backend", 28, RGBColor(0x00, 0x96, 0x88), True),
    ("Python FastAPI", 20, LIGHT, False),
    ("SQLite + aiosqlite", 20, LIGHT, False),
    ("Rate limiting", 20, LIGHT, False),
    ("Atomic transactions", 20, LIGHT, False),
])
add_multi(s, 9, 2.2, 3.5, 3, [
    ("AI", 28, RGBColor(0xD9, 0x77, 0x57), True),
    ("Claude Haiku 4.5", 20, LIGHT, False),
    ("3 dedicated prompts", 20, LIGHT, False),
    ("OpenAI Whisper", 20, LIGHT, False),
    ("Multi-provider fallback", 20, LIGHT, False),
])
add_text(s, 1, 5.5, 11, 0.5, "Built with Kiro IDE  •  Open Source  •  Code Reviewed  •  50+ Commits", 20, ACCENT)

# ─── Slide 9: Social Impact ───
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.8, 11, 1, "Social Impact", 48, WHITE, bold=True)
add_multi(s, 1, 2.2, 6, 4.5, [
    ("→  366M+ people affected by ADHD worldwide", 22, LIGHT, False),
    ("→  10% of American college students", 22, LIGHT, False),
    ("→  Executive dysfunction ≠ laziness", 22, LIGHT, False),
    ("→  Quick wins first = behavioral activation therapy", 22, GREEN, False),
    ("→  Encouraging tone = combats shame spiral", 22, GREEN, False),
    ("→  Voice input = lowers barrier to entry", 22, GREEN, False),
    ("→  No time estimates = reduces anxiety", 22, GREEN, False),
])
add_multi(s, 8, 3, 4.5, 2, [
    ("💙", 48, WHITE, False),
    ('"Executive dysfunction isn\'t laziness — it\'s a neurological difficulty with task initiation and planning."', 18, MUTED, False),
])

# ─── Slide 10: End ───
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 0, 1.5, 13.333, 1, "✨", 72, WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0, 2.5, 13.333, 1.2, "Clarity", 64, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 3.8, 10, 1, "You don't need to have it all figured out.\nJust tell us what's on your mind.", 28, MUTED, align=PP_ALIGN.CENTER)
add_text(s, 0, 5.2, 13.333, 0.5, "github.com/SeaL773/clarity", 22, ACCENT, align=PP_ALIGN.CENTER)
add_text(s, 0, 6.2, 13.333, 0.5, "Thank you.", 24, WHITE, align=PP_ALIGN.CENTER)

out = r"E:\github\hackathon\AWS Kiro x CS Careers Hackathon at Virginia Tech\hackathon\Clarity_Pitch.pptx"
prs.save(out)
print(f"Saved: {out}")
